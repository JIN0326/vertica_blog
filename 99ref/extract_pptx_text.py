from pathlib import Path
from pptx import Presentation

p = Path('c:/Users/SOO/MyProject/vertica_blog/Vertica_소개자료_202509.pptx')
print('exists', p.exists(), 'size', p.stat().st_size if p.exists() else None)
prs = Presentation(str(p))
for i, slide in enumerate(prs.slides, 1):
    print(f'--- Slide {i} ---')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            if text:
                print(text)
